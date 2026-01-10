#!/usr/bin/env python3
"""
ABOUTME: Converts PPTX slides to PNG images using python-pptx and PIL
ABOUTME: Renders slides with proper background colors and text positioning
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

def emu_to_px(emu, dpi=150):
    """Convert EMU (English Metric Units) to pixels at given DPI"""
    inches = emu / 914400
    return int(inches * dpi)

def hex_to_rgb(hex_color):
    """Convert hex color string to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def get_rgb_from_color(color):
    """Extract RGB values from pptx color object"""
    try:
        if hasattr(color, 'rgb') and color.rgb:
            rgb_str = str(color.rgb)
            if rgb_str and len(rgb_str) == 6:
                return hex_to_rgb(rgb_str)
    except Exception as e:
        pass
    return None

def render_slide_to_image(slide, prs, output_path, slide_num, dpi=150):
    """Render a single slide to an image"""
    width_px = emu_to_px(prs.slide_width, dpi)
    height_px = emu_to_px(prs.slide_height, dpi)

    # Default background - check slide background first
    bg_color = (255, 255, 255)  # White default

    # Try to get slide background color from slide's fill
    try:
        bg = slide.background
        if bg and bg.fill:
            fill = bg.fill
            if fill.type is not None:
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    rgb = get_rgb_from_color(fill.fore_color)
                    if rgb:
                        bg_color = rgb
    except Exception as e:
        pass

    # Create image with background color
    img = Image.new('RGB', (width_px, height_px), bg_color)
    draw = ImageDraw.Draw(img)

    # Try to load fonts
    try:
        font_regular = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        font_bold = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
        font_mono = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf", 18)
    except:
        font_regular = ImageFont.load_default()
        font_bold = font_regular
        font_mono = font_regular

    # Render shapes
    for shape in slide.shapes:
        try:
            x = emu_to_px(shape.left, dpi)
            y = emu_to_px(shape.top, dpi)
            w = emu_to_px(shape.width, dpi)
            h = emu_to_px(shape.height, dpi)

            # Draw shape fill
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                fill_color = get_rgb_from_color(shape.fill.fore_color)
                if fill_color:
                    # Check for outline
                    outline = None
                    if hasattr(shape, 'line') and shape.line and shape.line.color:
                        outline = get_rgb_from_color(shape.line.color)
                    draw.rectangle([x, y, x+w, y+h], fill=fill_color, outline=outline)

            # Draw text
            if shape.has_text_frame:
                text_y = y + 10
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    if not text:
                        continue

                    # Get font properties
                    font_size = 24
                    is_bold = False
                    text_color = (0, 0, 0)

                    if paragraph.runs:
                        run = paragraph.runs[0]
                        if run.font:
                            if run.font.size:
                                font_size = int(run.font.size.pt * dpi / 72)
                            if run.font.bold:
                                is_bold = True
                            if run.font.color and run.font.color.rgb:
                                rgb = get_rgb_from_color(run.font.color)
                                if rgb:
                                    text_color = rgb

                    # Select font based on properties
                    try:
                        if is_bold:
                            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", font_size)
                        else:
                            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size)
                    except:
                        font = font_regular

                    # Handle alignment
                    text_x = x + 10
                    if paragraph.alignment:
                        from pptx.enum.text import PP_ALIGN
                        if paragraph.alignment == PP_ALIGN.CENTER:
                            text_x = x + (w - len(text) * font_size * 0.5) / 2
                        elif paragraph.alignment == PP_ALIGN.RIGHT:
                            text_x = x + w - len(text) * font_size * 0.5 - 10

                    draw.text((text_x, text_y), text, fill=text_color, font=font)
                    text_y += font_size + 5

        except Exception as e:
            pass

    img.save(output_path)
    return output_path

def convert_pptx_to_images(pptx_path, output_dir, prefix):
    """Convert all slides in a PPTX to images"""
    print(f"Converting {pptx_path}...")

    prs = Presentation(pptx_path)
    os.makedirs(output_dir, exist_ok=True)

    for i, slide in enumerate(prs.slides, start=1):
        output_path = os.path.join(output_dir, f"{prefix}-slide-{i}.png")
        try:
            render_slide_to_image(slide, prs, output_path, i)
            print(f"  Created: {prefix}-slide-{i}.png")
        except Exception as e:
            print(f"  Failed slide {i}: {e}")

    return len(prs.slides)

if __name__ == "__main__":
    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    input_pptx = os.path.join(project_root, "output", "input-presentation.pptx")
    output_pptx = os.path.join(project_root, "output", "output-presentation.pptx")
    comparison_dir = os.path.join(project_root, "comparison")

    if os.path.exists(input_pptx):
        convert_pptx_to_images(input_pptx, comparison_dir, "input")

    if os.path.exists(output_pptx):
        convert_pptx_to_images(output_pptx, comparison_dir, "output")

    print("\nConversion complete!")
